"""
Yigthinker — PPTX rendering engine (ported from Yigcore skill-pptx, verdict ACCEPT).

This package bundles the pure-Python engine originally shipped in
`Yigcore/packages/skills/skill-pptx/engine/` — see
`docs/audit/2026-04-skill-pptx-audit.md` at commit 443ead3 for the audit
that approved the port.

Ported modules (kept for future richer use-cases):
- schema.py              — Pydantic models for template/slide inputs
- slide_creator.py       — `create_from_scratch` + 4 built-in themes
- chart_handler.py       — native PowerPoint chart insertion
- template_parser.py     — extract TemplateSchema from existing .pptx
- template_updater.py    — placeholder-aware content updates
- comment_writer.py      — OOXML review comments via lxml
- snapshot_generator.py  — LibreOffice PNG previews (lazy import;
                           no import-time side effects)

Public API for `report_generate` (Phase 1b minimum):

    render_pptx(df, output_path, title) -> Path

renders a simple deck with a title slide and a data table slide directly
via python-pptx. This is a thin inline orchestration — we deliberately
do NOT route the DataFrame case through `slide_creator.create_from_scratch`
because that API is template-oriented. The richer ported modules remain
available for future features (chart insertion, comment markup, etc.).
"""
from __future__ import annotations

from pathlib import Path

import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def _add_title_slide(prs: Presentation, title: str, subtitle: str | None = None) -> None:
    """Add a centered title slide using the default python-pptx title layout."""
    layout = prs.slide_layouts[0]  # "Title Slide"
    slide = prs.slides.add_slide(layout)
    if slide.placeholders and slide.placeholders[0].has_text_frame:
        slide.placeholders[0].text = title
        # Style the title a touch.
        for para in slide.placeholders[0].text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(40)
                run.font.bold = True
                run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)
    # Subtitle placeholder (idx=1) if present.
    if subtitle:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1 and ph.has_text_frame:
                ph.text = subtitle
                break


def _add_table_slide(prs: Presentation, df: pd.DataFrame, title: str) -> None:
    """Add a title + table slide rendering all rows of df (header row + data)."""
    # "Title Only" layout keeps the slide clean; fall back to Blank if unavailable.
    try:
        layout = prs.slide_layouts[5]  # Title Only
    except IndexError:
        layout = prs.slide_layouts[-1]
    slide = prs.slides.add_slide(layout)

    # Set title on the title placeholder if present.
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0 and ph.has_text_frame:
            ph.text = f"{title} — Data"
            break

    # Table dimensions.
    n_cols = max(len(df.columns), 1)
    n_rows = len(df) + 1  # header + data
    # Cap rows rendered on a single slide to keep output sensible; the full
    # DataFrame remains in ctx.vars for richer downstream use.
    MAX_ROWS = 40
    if n_rows > MAX_ROWS + 1:
        n_rows = MAX_ROWS + 1
    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(12.33)
    height = Inches(5.5)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Header row.
    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True
                run.font.size = Pt(12)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Data rows (bounded by MAX_ROWS).
    for row_idx, row in enumerate(df.head(MAX_ROWS).itertuples(index=False), start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = "" if value is None else str(value)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(10)


def _add_chart_slide_if_numeric(prs: Presentation, df: pd.DataFrame, title: str) -> None:
    """Append a bar-chart slide if df has at least one numeric column.

    Uses the first non-numeric column (or index as fallback) as categories
    and every numeric column as a series. Silently skips on any failure —
    chart slides are a bonus, not a guaranteed part of the output.
    """
    numeric_cols = [c for c in df.columns if pd.api.types.is_numeric_dtype(df[c])]
    if not numeric_cols:
        return
    try:
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        cat_col = next((c for c in df.columns if c not in numeric_cols), None)
        categories = (
            [str(v) for v in df[cat_col].tolist()]
            if cat_col is not None
            else [str(i) for i in range(len(df))]
        )

        # Cap categories to keep chart legible.
        MAX_CATS = 20
        categories = categories[:MAX_CATS]

        chart_data = CategoryChartData()
        chart_data.categories = categories
        for col in numeric_cols:
            values = [float(v) if pd.notna(v) else 0.0 for v in df[col].tolist()[:MAX_CATS]]
            chart_data.add_series(str(col), values)

        try:
            layout = prs.slide_layouts[5]  # Title Only
        except IndexError:
            layout = prs.slide_layouts[-1]
        slide = prs.slides.add_slide(layout)
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0 and ph.has_text_frame:
                ph.text = f"{title} — Chart"
                break

        left = Inches(1.0)
        top = Inches(1.3)
        width = Inches(11.33)
        height = Inches(5.5)
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
        )
    except Exception:
        # Non-fatal: render the deck without the chart slide.
        return


def render_pptx(
    df: pd.DataFrame,
    output_path: Path | str,
    title: str = "Report",
) -> Path:
    """Render a DataFrame to a .pptx deck.

    Produces:
    1. Title slide with `title`
    2. Table slide with df header + data rows (capped at 40 rows)
    3. Optional bar-chart slide if df has at least one numeric column

    The caller is responsible for path-safety checks; `report_generate`
    funnels `output_path` through `_safe_output_path` before invoking us.

    Returns the absolute Path of the saved .pptx.
    """
    prs = Presentation()
    # Standard widescreen 16:9.
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, title=title)
    _add_table_slide(prs, df, title=title)
    _add_chart_slide_if_numeric(prs, df, title=title)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out.resolve()


__all__ = ["render_pptx"]
