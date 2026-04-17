"""
skill-pptx — Template Parser

Parses a .pptx file and extracts its complete TemplateSchema:
- Slide masters and layouts (with placeholders)
- Each slide's content and structure
- Theme colors and fonts
"""

from typing import Optional
from pptx import Presentation
from pptx.util import Emu

from .schema import (
    PlaceholderInfo,
    SlideLayoutInfo,
    SlideInfo,
    ThemeInfo,
    TemplateSchema,
)


# Map python-pptx placeholder type enum values to readable strings
PP_PLACEHOLDER_TYPE_NAMES = {
    1: "title",
    2: "body",
    3: "center_title",
    4: "subtitle",
    5: "date",
    6: "slide_number",
    7: "footer",
    8: "object",
    9: "resource",
    10: "picture",
    11: "clipart",
    12: "diagram",
    13: "chart",
    14: "table",
    15: "header",
    16: "media_clip",
}


def _get_placeholder_type(ph) -> str:
    """Convert python-pptx placeholder type to readable string."""
    ph_type = ph.placeholder_format.type
    if ph_type is not None:
        type_val = ph_type.real if hasattr(ph_type, 'real') else int(ph_type)
        return PP_PLACEHOLDER_TYPE_NAMES.get(type_val, f"type_{type_val}")
    return "unknown"


def _extract_font_info(run) -> Optional[dict]:
    """Extract font info from a run object."""
    try:
        font = run.font
        color = None
        if font.color and font.color.type is not None:
            try:
                color = str(font.color.rgb)
            except Exception:
                color = None
        return {
            "name": font.name,
            "size": int(font.size) if font.size else None,
            "bold": font.bold,
            "color": color,
        }
    except Exception:
        return None


def _extract_placeholder(ph, extract_text: bool = False) -> PlaceholderInfo:
    """Extract PlaceholderInfo from a placeholder shape."""
    position = {
        "left": ph.left or 0,
        "top": ph.top or 0,
        "width": ph.width or 0,
        "height": ph.height or 0,
    }

    current_text = None
    font_info = None

    if extract_text and ph.has_text_frame:
        current_text = ph.text_frame.text
        # Get font info from first run of first paragraph
        try:
            first_para = ph.text_frame.paragraphs[0]
            if first_para.runs:
                font_info = _extract_font_info(first_para.runs[0])
        except (IndexError, AttributeError):
            pass

    return PlaceholderInfo(
        idx=ph.placeholder_format.idx,
        type=_get_placeholder_type(ph),
        name=ph.name,
        position=position,
        has_text=ph.has_text_frame,
        current_text=current_text,
        font_info=font_info,
    )


def _find_layout_index(prs: Presentation, target_layout) -> int:
    """Find the index of a layout within all masters."""
    idx = 0
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout == target_layout:
                return idx
            idx += 1
    return -1


def _extract_theme(prs: Presentation) -> ThemeInfo:
    """Extract theme colors and fonts from the first slide master."""
    colors = {}
    fonts = {}
    warnings = []

    try:
        master = prs.slide_masters[0]
        theme_element = master.element.find(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}theme"
        )

        if theme_element is not None:
            ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
            # Extract color scheme
            clr_scheme = theme_element.find(f".//{{{ns}}}clrScheme")
            if clr_scheme is not None:
                color_names = [
                    "dk1", "lt1", "dk2", "lt2",
                    "accent1", "accent2", "accent3",
                    "accent4", "accent5", "accent6",
                    "hlink", "folHlink"
                ]
                for name in color_names:
                    elem = clr_scheme.find(f"{{{ns}}}{name}")
                    if elem is not None:
                        # Try srgbClr first, then sysClr
                        srgb = elem.find(f"{{{ns}}}srgbClr")
                        if srgb is not None:
                            colors[name] = srgb.get("val", "")
                        else:
                            sys_clr = elem.find(f"{{{ns}}}sysClr")
                            if sys_clr is not None:
                                colors[name] = sys_clr.get("lastClr", "")

            # Extract font scheme
            font_scheme = theme_element.find(f".//{{{ns}}}fontScheme")
            if font_scheme is not None:
                major = font_scheme.find(f".//{{{ns}}}majorFont/{{{ns}}}latin")
                minor = font_scheme.find(f".//{{{ns}}}minorFont/{{{ns}}}latin")
                if major is not None:
                    fonts["major"] = major.get("typeface", "")
                if minor is not None:
                    fonts["minor"] = minor.get("typeface", "")
    except Exception as e:
        warnings.append(f"Could not extract theme info: {str(e)}")

    return ThemeInfo(colors=colors, fonts=fonts)


def parse_template(file_path: str) -> dict:
    """
    Parse a .pptx file and return a TemplateSchema dict.

    Args:
        file_path: Path to the .pptx file

    Returns:
        dict matching TemplateSchema model
    """
    prs = Presentation(file_path)
    warnings = []

    schema = TemplateSchema(
        slide_count=len(prs.slides),
        slide_width=prs.slide_width,
        slide_height=prs.slide_height,
        layouts=[],
        slides=[],
        theme=_extract_theme(prs),
        warnings=warnings,
    )

    # Extract all available layouts from all masters
    for master in prs.slide_masters:
        master_name = master.name or "Default"
        for i, layout in enumerate(master.slide_layouts):
            placeholders = []
            for ph in layout.placeholders:
                try:
                    placeholders.append(_extract_placeholder(ph))
                except Exception as e:
                    warnings.append(f"Could not parse placeholder in layout '{layout.name}': {e}")

            layout_info = SlideLayoutInfo(
                name=layout.name,
                index=i,
                placeholders=placeholders,
                master_name=master_name,
            )
            schema.layouts.append(layout_info)

    # Extract each existing slide's content
    for idx, slide in enumerate(prs.slides):
        placeholders = []
        for ph in slide.placeholders:
            try:
                placeholders.append(_extract_placeholder(ph, extract_text=True))
            except Exception as e:
                warnings.append(f"Could not parse placeholder on slide {idx}: {e}")

        # Detect unsupported features
        has_chart = False
        has_table = False
        for shape in slide.shapes:
            if hasattr(shape, 'has_chart') and shape.has_chart:
                has_chart = True
            if hasattr(shape, 'has_table') and shape.has_table:
                has_table = True
            # Warn about complex unsupported elements
            shape_type = shape.shape_type
            # 13 = MSO_SHAPE_TYPE.MEDIA  (videos, audio)
            if shape_type == 16:  # LINKED_OLE_OBJECT
                warnings.append(f"Slide {idx} contains linked OLE object (SmartArt/3D) which may not be fully supported")

        notes_text = None
        if slide.has_notes_slide:
            try:
                notes_text = slide.notes_slide.notes_text_frame.text
            except Exception:
                pass

        slide_info = SlideInfo(
            index=idx,
            layout_name=slide.slide_layout.name,
            layout_index=_find_layout_index(prs, slide.slide_layout),
            placeholders=placeholders,
            shapes_count=len(slide.shapes),
            has_chart=has_chart,
            has_table=has_table,
            has_notes=slide.has_notes_slide,
            notes_text=notes_text,
        )
        schema.slides.append(slide_info)

    return schema.model_dump()


def analyze_content(file_path: str, slide_indices: Optional[list] = None) -> dict:
    """
    Analyze existing slide content in detail.

    Args:
        file_path: Path to the .pptx file
        slide_indices: Optional list of 0-based slide indices to analyze

    Returns:
        dict with detailed content analysis per slide
    """
    prs = Presentation(file_path)
    results = []

    indices_to_analyze = slide_indices if slide_indices is not None else list(range(len(prs.slides)))

    for idx in indices_to_analyze:
        if idx >= len(prs.slides):
            continue

        slide = prs.slides[idx]
        slide_data = {
            "slide_index": idx,
            "layout": slide.slide_layout.name,
            "placeholders": [],
            "shapes": [],
            "has_notes": slide.has_notes_slide,
            "notes_text": None,
        }

        # Extract placeholder content
        for ph in slide.placeholders:
            ph_data = {
                "idx": ph.placeholder_format.idx,
                "type": _get_placeholder_type(ph),
                "name": ph.name,
                "text": None,
                "paragraphs": [],
            }

            if ph.has_text_frame:
                ph_data["text"] = ph.text_frame.text
                for para in ph.text_frame.paragraphs:
                    para_data = {
                        "text": para.text,
                        "runs": [],
                    }
                    for run in para.runs:
                        font_info = _extract_font_info(run)
                        para_data["runs"].append({
                            "text": run.text,
                            "font": font_info,
                        })
                    ph_data["paragraphs"].append(para_data)

            # Extract chart data summary
            if hasattr(ph, 'has_chart') and ph.has_chart:
                try:
                    chart = ph.chart
                    ph_data["chart_type"] = str(chart.chart_type)
                    ph_data["chart_series_count"] = len(chart.series)
                except Exception:
                    pass

            # Extract table data summary
            if hasattr(ph, 'has_table') and ph.has_table:
                try:
                    table = ph.table
                    ph_data["table_rows"] = table.rows.__len__()
                    ph_data["table_cols"] = table.columns.__len__()
                except Exception:
                    pass

            slide_data["placeholders"].append(ph_data)

        # Extract non-placeholder shapes summary
        for shape in slide.shapes:
            if not shape.is_placeholder:
                shape_data = {
                    "name": shape.name,
                    "shape_type": str(shape.shape_type),
                    "has_text": shape.has_text_frame,
                }
                if shape.has_text_frame:
                    shape_data["text"] = shape.text_frame.text
                slide_data["shapes"].append(shape_data)

        # Extract notes
        if slide.has_notes_slide:
            try:
                slide_data["notes_text"] = slide.notes_slide.notes_text_frame.text
            except Exception:
                pass

        results.append(slide_data)

    return {
        "file_path": file_path,
        "slide_count": len(prs.slides),
        "analyzed_slides": results,
    }
