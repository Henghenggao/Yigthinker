"""
skill-pptx — Template Updater

Updates slide content based on a ContentUpdatePlan.
Reads existing .pptx, updates placeholders preserving run-level formatting.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

import copy
import os


def _update_text_preserving_format(placeholder, new_text: str):
    """
    Run-level text replacement that preserves original font/color/size.

    Strategy:
    1. Save formatting from first run of first paragraph
    2. Clear all but first paragraph
    3. Clear all but first run in first paragraph
    4. Set new text on first run
    5. Restore saved formatting
    """
    if not placeholder.has_text_frame:
        return

    text_frame = placeholder.text_frame
    paragraphs = text_frame.paragraphs

    if not paragraphs:
        text_frame.text = new_text
        return

    first_para = paragraphs[0]
    runs = first_para.runs

    if runs:
        first_run = runs[0]
        # Save formatting
        saved_font_name = first_run.font.name
        saved_font_size = first_run.font.size
        saved_font_bold = first_run.font.bold
        saved_font_italic = first_run.font.italic
        saved_color = None
        try:
            if first_run.font.color and first_run.font.color.type is not None:
                saved_color = first_run.font.color.rgb
        except Exception:
            pass

        # Clear extra runs in first paragraph (keep only first)
        for run in runs[1:]:
            run.text = ""

        # Clear extra paragraphs (keep only first)
        # We can't directly delete paragraphs in python-pptx easily,
        # so we just clear their content
        for para in paragraphs[1:]:
            for run in para.runs:
                run.text = ""

        # Set new text on first run
        first_run.text = new_text

        # Restore formatting
        if saved_font_name is not None:
            first_run.font.name = saved_font_name
        if saved_font_size is not None:
            first_run.font.size = saved_font_size
        if saved_font_bold is not None:
            first_run.font.bold = saved_font_bold
        if saved_font_italic is not None:
            first_run.font.italic = saved_font_italic
        if saved_color is not None:
            try:
                from pptx.dml.color import RGBColor
                first_run.font.color.rgb = RGBColor.from_string(str(saved_color))
            except Exception:
                pass
    else:
        # No existing runs — just set text directly
        first_para.clear()
        run = first_para.add_run()
        run.text = new_text


def _update_chart_data(placeholder, chart_data: dict):
    """Update chart data for a chart placeholder."""
    if not (hasattr(placeholder, 'has_chart') and placeholder.has_chart):
        return

    try:
        chart = placeholder.chart
        new_chart_data = CategoryChartData()

        categories = chart_data.get("categories", [])
        new_chart_data.categories = categories

        for series in chart_data.get("series", []):
            new_chart_data.add_series(
                series.get("name", "Series"),
                series.get("values", [])
            )

        chart.replace_data(new_chart_data)
    except Exception as e:
        raise RuntimeError(f"Failed to update chart data: {e}")


def _update_table_data(placeholder, table_data: dict):
    """Update table content for a table placeholder."""
    if not (hasattr(placeholder, 'has_table') and placeholder.has_table):
        return

    try:
        table = placeholder.table
        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])

        # Update header row (row 0)
        if headers and table.rows:
            for col_idx, header_text in enumerate(headers):
                if col_idx < len(table.columns):
                    cell = table.cell(0, col_idx)
                    cell.text = str(header_text)

        # Update data rows
        for row_idx, row_data in enumerate(rows):
            actual_row = row_idx + 1  # Skip header
            if actual_row < len(table.rows):
                for col_idx, cell_value in enumerate(row_data):
                    if col_idx < len(table.columns):
                        cell = table.cell(actual_row, col_idx)
                        cell.text = str(cell_value) if cell_value is not None else ""
    except Exception as e:
        raise RuntimeError(f"Failed to update table data: {e}")


def _find_placeholder_by_idx(slide, ph_idx: int):
    """Find a placeholder by its idx (not its position index)."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == ph_idx:
            return ph
    return None


def update_slides(template_path: str, updates: list, output_path: str) -> dict:
    """
    Update slide content based on a ContentUpdatePlan.

    Args:
        template_path: Path to the template .pptx
        updates: List of slide update dicts [{slideIndex, placeholderUpdates}]
        output_path: Output file path

    Returns:
        dict with operation result
    """
    prs = Presentation(template_path)
    slides_updated = 0
    placeholders_updated = 0
    errors = []

    for update in updates:
        slide_idx = update.get("slideIndex", update.get("slide_index"))
        if slide_idx is None or slide_idx >= len(prs.slides):
            errors.append(f"Slide index {slide_idx} out of range (total: {len(prs.slides)})")
            continue

        slide = prs.slides[slide_idx]
        ph_updates = update.get("placeholderUpdates", update.get("placeholder_updates", []))

        for ph_update in ph_updates:
            ph_idx = ph_update.get("placeholderIdx", ph_update.get("placeholder_idx"))
            update_type = ph_update.get("type")
            content = ph_update.get("content")

            placeholder = _find_placeholder_by_idx(slide, ph_idx)
            if placeholder is None:
                errors.append(f"Placeholder idx={ph_idx} not found on slide {slide_idx}")
                continue

            try:
                if update_type == "text":
                    _update_text_preserving_format(placeholder, str(content))
                    placeholders_updated += 1
                elif update_type == "chart_data":
                    _update_chart_data(placeholder, content if isinstance(content, dict) else {})
                    placeholders_updated += 1
                elif update_type == "image":
                    if isinstance(content, str) and os.path.exists(content):
                        placeholder.insert_picture(content)
                        placeholders_updated += 1
                    else:
                        errors.append(f"Image path not found: {content}")
                elif update_type == "table":
                    _update_table_data(placeholder, content if isinstance(content, dict) else {})
                    placeholders_updated += 1
                else:
                    errors.append(f"Unknown update type: {update_type}")
            except Exception as e:
                errors.append(f"Error updating slide {slide_idx}, placeholder {ph_idx}: {str(e)}")

        slides_updated += 1

    # Ensure output directory exists
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)

    return {
        "success": True,
        "output_path": output_path,
        "slides_updated": slides_updated,
        "placeholders_updated": placeholders_updated,
        "errors": errors,
    }
