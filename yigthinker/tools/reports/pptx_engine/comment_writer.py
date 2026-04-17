"""
skill-pptx — Comment Writer

Adds OOXML review comments to PowerPoint slides via lxml.
Comments appear in PowerPoint's review panel.

Reference: ECMA-376 Part 1, §19.4 (PresentationML Comments)
"""

import os
from datetime import datetime
from lxml import etree
from pptx import Presentation
from pptx.opc.packuri import PackURI
from pptx.opc.package import Part


# OOXML namespaces
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
CM_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Relationship types
REL_TYPE_COMMENT_AUTHORS = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/"
    "commentAuthors"
)
REL_TYPE_COMMENTS = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments"
)

# Content types
CT_COMMENT_AUTHORS = (
    "application/vnd.openxmlformats-officedocument.presentationml."
    "commentAuthors+xml"
)
CT_COMMENTS = (
    "application/vnd.openxmlformats-officedocument.presentationml.comments+xml"
)


def _get_or_create_comment_authors(prs: Presentation) -> tuple:
    """
    Ensure commentAuthors.xml exists in the package.

    Returns:
        (authors_element, author_map dict{name: id})
    """
    # Build the commentAuthors XML
    nsmap = {"p": PML_NS}
    authors_root = etree.Element(f"{{{PML_NS}}}cmAuthorLst", nsmap=nsmap)
    author_map = {}

    # Check if commentAuthors part already exists via presentation relationships
    try:
        prs_part = prs.part
        for rel in prs_part.rels.values():
            if rel.reltype == REL_TYPE_COMMENT_AUTHORS:
                # Parse existing authors
                existing_xml = rel.target_part._blob
                existing_root = etree.fromstring(existing_xml)
                for author_elem in existing_root.findall(f"{{{PML_NS}}}cmAuthor"):
                    name = author_elem.get("name", "")
                    author_id = author_elem.get("id", "0")
                    author_map[name] = int(author_id)
                return existing_root, author_map
    except Exception:
        pass

    return authors_root, author_map


def _ensure_author(authors_root, author_map: dict, author_name: str) -> int:
    """Get or create an author in the commentAuthors element."""
    if author_name in author_map:
        return author_map[author_name]

    new_id = max(author_map.values(), default=-1) + 1
    author_map[author_name] = new_id

    author_elem = etree.SubElement(
        authors_root,
        f"{{{PML_NS}}}cmAuthor",
    )
    author_elem.set("id", str(new_id))
    author_elem.set("name", author_name)
    author_elem.set("initials", author_name[:2].upper())
    author_elem.set("lastIdx", "0")
    author_elem.set("clrIdx", str(new_id % 8))

    return new_id


def _build_comment_element(
    author_id: int,
    comment_idx: int,
    text: str,
    pos_x: int = 1270000,
    pos_y: int = 1270000,
) -> etree._Element:
    """Build a <p:cm> element for a single comment."""
    dt = datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%S")

    cm = etree.Element(f"{{{PML_NS}}}cm")
    cm.set("authorId", str(author_id))
    cm.set("dt", dt)
    cm.set("idx", str(comment_idx))

    # Position
    pos = etree.SubElement(cm, f"{{{PML_NS}}}pos")
    pos.set("x", str(pos_x))
    pos.set("y", str(pos_y))

    # Text
    text_elem = etree.SubElement(cm, f"{{{PML_NS}}}text")
    text_elem.text = text

    return cm


def add_comments(file_path: str, comments: list, output_path: str) -> dict:
    """
    Add review comments to slides.

    Args:
        file_path: Path to the .pptx file
        comments: List of comment dicts [{slideIndex, author, text, position?}]
        output_path: Output file path

    Returns:
        dict with operation result
    """
    prs = Presentation(file_path)
    added = 0
    errors = []

    # Group comments by slide
    by_slide: dict = {}
    for comment in comments:
        slide_idx = comment.get("slideIndex", comment.get("slide_index", 0))
        if slide_idx not in by_slide:
            by_slide[slide_idx] = []
        by_slide[slide_idx].append(comment)

    # Build authors XML
    nsmap = {"p": PML_NS}
    authors_root = etree.Element(f"{{{PML_NS}}}cmAuthorLst", nsmap=nsmap)
    author_map: dict = {}

    for slide_idx, slide_comments in by_slide.items():
        if slide_idx >= len(prs.slides):
            errors.append(f"Slide index {slide_idx} out of range")
            continue

        slide = prs.slides[slide_idx]

        # Build comments XML for this slide
        cm_lst = etree.Element(f"{{{PML_NS}}}cmLst", nsmap={"p": PML_NS})
        comment_idx = 1

        for comment in slide_comments:
            author_name = comment.get("author", "AI Assistant")
            text = comment.get("text", "")
            position = comment.get("position", {})
            pos_x = position.get("x", 1270000) if position else 1270000
            pos_y = position.get("y", 1270000) if position else 1270000

            author_id = _ensure_author(authors_root, author_map, author_name)

            cm_elem = _build_comment_element(
                author_id=author_id,
                comment_idx=comment_idx,
                text=text,
                pos_x=pos_x,
                pos_y=pos_y,
            )
            cm_lst.append(cm_elem)
            comment_idx += 1
            added += 1

        # Add comments part to slide
        try:
            # Serialize comments XML
            cm_xml = etree.tostring(cm_lst, xml_declaration=True, encoding="UTF-8", standalone=True)

            # Create a new part for slide comments
            slide_part = slide.part
            comments_uri = PackURI(f"/ppt/comments/comment{slide_idx + 1}.xml")

            comments_part = Part(
                comments_uri,
                CT_COMMENTS,
                slide_part.package,
                cm_xml,
            )

            # Add relationship from slide to comments
            slide_part.relate_to(comments_part, REL_TYPE_COMMENTS)
        except Exception as e:
            errors.append(f"Failed to add comments to slide {slide_idx}: {str(e)}")

    # Save commentAuthors if we added any authors
    if author_map:
        try:
            authors_xml = etree.tostring(
                authors_root, xml_declaration=True, encoding="UTF-8", standalone=True
            )
            authors_uri = PackURI("/ppt/commentAuthors.xml")
            prs_part = prs.part

            authors_part = Part(
                authors_uri,
                CT_COMMENT_AUTHORS,
                prs_part.package,
                authors_xml,
            )
            prs_part.relate_to(authors_part, REL_TYPE_COMMENT_AUTHORS)
        except Exception as e:
            errors.append(f"Failed to save commentAuthors: {str(e)}")

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)

    return {
        "success": True,
        "comments_added": added,
        "output_path": output_path,
        "errors": errors,
    }
