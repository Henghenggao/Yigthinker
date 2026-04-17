"""
skill-pptx — Chart Handler

Inserts native PowerPoint charts (editable, not images) into slides.
Uses python-pptx's chart API for full OOXML compliance.
"""

import os
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE


# Map string chart type to python-pptx enum
CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "area": XL_CHART_TYPE.AREA,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "waterfall": XL_CHART_TYPE.COLUMN_CLUSTERED,  # python-pptx lacks native waterfall; use column
}


def insert_chart(
    file_path: str,
    slide_index: int,
    chart_type: str,
    chart_data: dict,
    title: str,
    position: dict,
    output_path: str,
) -> dict:
    """
    Insert a native PowerPoint chart into a slide.

    Args:
        file_path: Path to the .pptx file
        slide_index: 0-based slide index
        chart_type: "bar", "line", "pie", "area", "scatter", "waterfall"
        chart_data: { categories: [...], series: [{name, values}] }
        title: Chart title
        position: { left, top, width, height } in inches (optional)
        output_path: Output file path

    Returns:
        dict with operation result
    """
    prs = Presentation(file_path)

    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide index {slide_index} out of range (total: {len(prs.slides)})")

    slide = prs.slides[slide_index]

    # Determine chart type
    xl_chart_type = CHART_TYPE_MAP.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    # Build chart data
    if chart_type == "scatter":
        cd = XyChartData()
        for series in chart_data.get("series", []):
            series_data = cd.add_series(series.get("name", "Series"))
            values = series.get("values", [])
            categories = chart_data.get("categories", [])
            for i, val in enumerate(values):
                x = i
                try:
                    x = float(categories[i]) if i < len(categories) else i
                except (ValueError, TypeError):
                    x = i
                series_data.add_data_point(x, val)
    else:
        cd = CategoryChartData()
        cd.categories = chart_data.get("categories", [])
        for series in chart_data.get("series", []):
            cd.add_series(series.get("name", "Series"), series.get("values", []))

    # Determine position
    if position:
        left = Inches(position.get("left", 1.0))
        top = Inches(position.get("top", 1.5))
        width = Inches(position.get("width", 8.0))
        height = Inches(position.get("height", 4.5))
    else:
        left = Inches(1.0)
        top = Inches(1.5)
        width = Inches(8.0)
        height = Inches(4.5)

    # Add chart to slide
    graphic_frame = slide.shapes.add_chart(xl_chart_type, left, top, width, height, cd)
    chart = graphic_frame.chart

    # Set title
    if title:
        chart.has_title = True
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = title

    # Common chart formatting
    chart.has_legend = len(chart_data.get("series", [])) > 1

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)

    return {
        "success": True,
        "output_path": output_path,
        "chart_type": chart_type,
        "series_count": len(chart_data.get("series", [])),
        "category_count": len(chart_data.get("categories", [])),
    }
