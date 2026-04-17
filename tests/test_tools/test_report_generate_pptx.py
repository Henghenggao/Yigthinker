"""E2E lock: report_generate produces a real .pptx file (Phase 1b D2).

Ported engine lives at `yigthinker/tools/reports/pptx_engine/` — verdict
ACCEPT from `docs/audit/2026-04-skill-pptx-audit.md` (commit 443ead3).
"""
from __future__ import annotations

import pandas as pd
import pytest

from yigthinker.session import SessionContext
from yigthinker.tools.reports.report_generate import (
    ReportGenerateInput,
    ReportGenerateTool,
)


@pytest.mark.asyncio
async def test_report_generate_pptx_creates_file(tmp_path):
    pytest.importorskip("pptx")

    ctx = SessionContext(settings={"workspace_dir": str(tmp_path)})
    ctx.vars.set(
        "df1",
        pd.DataFrame({"metric": ["rev", "cost"], "value": [100, 40]}),
    )

    tool = ReportGenerateTool()
    out = tmp_path / "test.pptx"
    result = await tool.execute(
        ReportGenerateInput(
            var_name="df1",
            format="pptx",
            output_path=str(out),
            title="Test",
        ),
        ctx,
    )
    assert not result.is_error, result.content
    assert out.exists()
    assert out.stat().st_size > 1024  # non-trivial file

    # Magic bytes — .pptx is a zip archive.
    with out.open("rb") as f:
        assert f.read(2) == b"PK"

    # Smoke-open with python-pptx to make sure it's a valid deck.
    from pptx import Presentation

    prs = Presentation(str(out))
    # Expect >=2 slides: title + table (chart slide appears only if numeric col).
    assert len(prs.slides) >= 2


@pytest.mark.asyncio
async def test_report_generate_pptx_includes_chart_for_numeric_df(tmp_path):
    """With a numeric column, we expect an extra chart slide."""
    pytest.importorskip("pptx")

    ctx = SessionContext(settings={"workspace_dir": str(tmp_path)})
    ctx.vars.set(
        "sales",
        pd.DataFrame(
            {
                "region": ["EU", "US", "APAC"],
                "revenue": [120.0, 95.5, 60.0],
            }
        ),
    )

    tool = ReportGenerateTool()
    out = tmp_path / "sales.pptx"
    result = await tool.execute(
        ReportGenerateInput(
            var_name="sales",
            format="pptx",
            output_path=str(out),
            title="Sales",
        ),
        ctx,
    )
    assert not result.is_error, result.content

    from pptx import Presentation

    prs = Presentation(str(out))
    # Title + table + chart.
    assert len(prs.slides) == 3
